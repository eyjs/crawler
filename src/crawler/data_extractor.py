# src/crawler/data_extractor.py

import logging
import asyncio
import aiohttp
from bs4 import BeautifulSoup
from urllib.parse import urljoin, urlparse, parse_qs, unquote
from pathlib import Path
import io
import tempfile
import subprocess
import json
import time
import re
from typing import Dict, List, Tuple, Callable

# --- 각 파일 타입을 처리하기 위한 라이브러리 임포트 ---
from pypdf import PdfReader
from docx import Document
import pandas as pd
from pptx import Presentation
import hwp5

logger = logging.getLogger(__name__)

class DataExtractor:
    """
    HTML 및 다양한 첨부 파일에서 데이터를 추출하는 지능형 추출기.
    보일러플레이트 제거, 노이즈 필터링, 파싱 실패 격리 로직을 포함한 최종 버전.
    """
    def __init__(self):
        self._session = aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=90))
        self.failed_attachments_dir = Path("failed_attachments")
        self.failed_attachments_dir.mkdir(exist_ok=True)

        self.file_parsers: Dict[str, Callable] = {
            '.pdf': self._parse_pdf, '.docx': self._parse_docx,
            '.xlsx': self._parse_excel, '.xls': self._parse_excel,
            '.pptx': self._parse_pptx, '.hwp': self._parse_hwp,
        }

    def _clean_text(self, text: str) -> str:
        """추출된 텍스트에서 일반적인 노이즈 패턴과 불필요한 공백을 제거합니다."""
        # 여러 줄바꿈을 최대 두 개로 압축
        text = re.sub(r'(\n\s*){3,}', '\n\n', text)
        lines = text.split('\n')
        cleaned_lines = []

        noise_keywords = [
            '다운로드', '뷰어', '첨부파일', '목록으로', '이전글', '다음글', '맨위로',
            'Copyright', 'All rights reserved', '찾아오시는 길', '개인정보처리방침'
        ]
        for line in lines:
            stripped_line = line.strip()
            if not stripped_line: continue
            if len(stripped_line) < 10 and not re.search('[가-힣a-zA-Z0-9]', stripped_line): continue
            if any(keyword in stripped_line for keyword in noise_keywords): continue
            if re.match(r'^\s*(작성자|등록일|조회수|담당부서|키워드|분류)\s*[:\s]', stripped_line): continue
            cleaned_lines.append(stripped_line)

        return "\n".join(cleaned_lines)

    # --- 각 파일 타입별 텍스트 추출 메서드 ---
    def _parse_pdf(self, content: bytes) -> str:
        return "\n".join([page.extract_text() for page in PdfReader(io.BytesIO(content)).pages if page.extract_text()])
    def _parse_docx(self, content: bytes) -> str:
        return "\n".join([para.text for para in Document(io.BytesIO(content)).paragraphs if para.text])
    def _parse_excel(self, content: bytes) -> str:
        xls = pd.ExcelFile(io.BytesIO(content))
        return "\n\n".join([xls.parse(sheet_name).to_string() for sheet_name in xls.sheet_names])
    def _parse_pptx(self, content: bytes) -> str:
        prs = Presentation(io.BytesIO(content))
        return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    def _parse_hwp(self, content: bytes) -> str:
        try:
            hwp = hwp5.HWPFile(io.BytesIO(content))
            text = hwp.body.text()
            if not text: raise ValueError("Extracted HWP text is empty.")
            return text
        except Exception as e:
            raise e # 예외를 발생시켜 _extract_text_from_attachment에서 처리하도록 함

    async def _save_failed_attachment(self, url: str, content_bytes: bytes, reason: str, site_identifier: str):
        try:
            site_dir = self.failed_attachments_dir / site_identifier
            site_dir.mkdir(exist_ok=True)
            timestamp = int(time.time())
            original_filename = Path(urlparse(url).path).name or f"{timestamp}.attachment"
            file_path = site_dir / f"{timestamp}_{original_filename}"
            with open(file_path, 'wb') as f: f.write(content_bytes)
            meta_path = site_dir / f"{timestamp}_{original_filename}.meta.json"
            meta_data = {
                "original_url": url, "saved_path": str(file_path), "failure_reason": reason,
                "timestamp": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime(timestamp))
            }
            with open(meta_path, 'w', encoding='utf-8') as f: json.dump(meta_data, f, ensure_ascii=False, indent=2)
            logger.warning(f"🚨 파싱 실패! 원본 파일 및 메타데이터 저장 완료: {file_path}")
        except Exception as e:
            logger.error(f"실패한 첨부 파일 저장 중 오류 발생: {e}")

    async def _extract_text_from_attachment(self, url: str, site_identifier: str) -> str:
        file_ext = Path(urlparse(url).path).suffix.lower()
        parser = self.file_parsers.get(file_ext)
        if not parser: return ""
        content_bytes = None
        try:
            async with self._session.get(url, timeout=60) as response:
                if response.status != 200: return ""
                content_bytes = await response.read()
            text = await asyncio.to_thread(parser, content_bytes)
            if not text: raise ValueError("파서가 빈 텍스트를 반환했습니다.")
            logger.info(f"✅ {file_ext.upper()} 텍스트 추출 성공: {url} ({len(text)}자)")
            return f"\n\n--- 첨부된 {file_ext.upper()} 파일 내용: {url} ---\n{text}"
        except Exception as e:
            logger.error(f"첨부 파일({url}) 처리 실패: {e}")
            if content_bytes:
                await self._save_failed_attachment(url, content_bytes, str(e), site_identifier)
            return f"\n\n--- 첨부된 {file_ext.upper()} 파일 내용: {url} ---"

    def _get_real_pdf_url_from_viewer(self, viewer_url: str, base_url: str) -> str | None:
        try:
            query_params = parse_qs(urlparse(viewer_url).query)
            if 'file' in query_params:
                return urljoin(base_url, unquote(query_params['file'][0]))
        except Exception: return None

    async def extract(self, url: str, base_url: str, site_identifier: str) -> Dict | None:
        try:
            async with self._session.get(url, timeout=30) as response:
                if response.status != 200 or 'text/html' not in response.headers.get('Content-Type', ''):
                    return None
                html = await response.text()
                soup = BeautifulSoup(html, 'html.parser')

                title = soup.title.string.strip() if soup.title else ""

                for tag in soup(['header', 'footer', 'nav', 'aside', 'script', 'style', 'form', 'button']):
                    tag.decompose()

                main_content = soup.find('main') or soup.find('article') or soup.find('div', id='content') or soup.find('div', class_='content') or soup.find('div', id='bodyContent') or soup.body

                raw_text = main_content.get_text(separator='\n', strip=False) if main_content else ""
                main_text = self._clean_text(raw_text)

                all_links, attachment_links = [], []
                base_netloc = urlparse(base_url).netloc

                for a_tag in soup.find_all('a', href=True):
                    href = a_tag['href']
                    absolute_url = urljoin(base_url, href)

                    real_pdf_url = self._get_real_pdf_url_from_viewer(absolute_url, base_url)
                    if real_pdf_url and real_pdf_url not in attachment_links:
                        attachment_links.append(real_pdf_url); continue

                    file_ext = Path(urlparse(href).path).suffix.lower()
                    if file_ext in self.file_parsers and urlparse(absolute_url).netloc == base_netloc and absolute_url not in attachment_links:
                        attachment_links.append(absolute_url)
                    else:
                        all_links.append((absolute_url, a_tag.get_text(strip=True)))

                if attachment_links:
                    tasks = [self._extract_text_from_attachment(link, site_identifier) for link in attachment_links]
                    main_text += "".join(await asyncio.gather(*tasks))

                return {"url": url, "title": title, "main_text": main_text, "links": all_links}
        except Exception as e:
            logger.error(f"데이터 추출 중 예외 발생: {url}, {e}", exc_info=True)
            return None

    async def close_session(self):
        if self._session and not self._session.closed: await self._session.close()