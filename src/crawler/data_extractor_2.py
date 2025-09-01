import logging
import asyncio
import aiohttp
from bs4 import BeautifulSoup, Tag
from urllib.parse import urljoin, urlparse, parse_qs, unquote
from pathlib import Path
import io
import tempfile
import subprocess
import json
import time
import re
from typing import Dict, List, Tuple, Callable, Optional

# --- 각 파일 타입을 처리하기 위한 라이브러리 임포트 ---
from pypdf import PdfReader
from docx import Document
import pandas as pd
from pptx import Presentation
import hwp5

logger = logging.getLogger(__name__)

class DataExtractor:
    """
    [v2.6] '분석'과 '제거' 로직을 명확히 분리하여 'NoneType' 오류를
    근본적으로 해결한 최종 안정화 버전.
    """
    def __init__(self):
        timeout_config = aiohttp.ClientTimeout(total=120, connect=30, sock_read=90)
        self._session = aiohttp.ClientSession(timeout=timeout_config)

        self.failed_attachments_dir = Path("failed_attachments")
        self.failed_attachments_dir.mkdir(exist_ok=True)

        self.file_parsers: Dict[str, Callable] = {
            '.pdf': self._parse_pdf, '.docx': self._parse_docx,
            '.xlsx': self._parse_excel, '.xls': self._parse_excel,
            '.pptx': self._parse_pptx, '.hwp': self._parse_hwp,
        }

    # --- 1. 지능형 본문 추출 (알고리즘 안정성 강화) ---

    def _calculate_content_score(self, element: Tag) -> float:
        """ 요소의 '내용 밀도' 점수를 계산합니다. """
        if not isinstance(element, Tag) or element.name in ['script', 'style', 'a']:
            return 0

        text = element.get_text(strip=True)
        text_length = len(text)

        if text_length < 100:
            return 0

        link_text_length = sum(len(a.get_text(strip=True)) for a in element.find_all('a'))
        pure_text_length = text_length - link_text_length
        link_density = link_text_length / max(1, text_length)
        score = pure_text_length * (1 - link_density**2)

        tag_id = element.get('id', '').lower()
        tag_class = ' '.join(element.get('class', [])).lower()

        if any(keyword in tag_id or keyword in tag_class for keyword in ['content', 'article', 'post', 'body', 'main', 'view']):
            score *= 1.5

        return score

    def _get_best_candidate(self, soup: BeautifulSoup) -> Optional[Tag]:
        """ 가장 안정적인 방식으로 본문 후보를 찾습니다. """
        best_candidate = None
        highest_score = 0

        for element in soup.find_all(['div', 'article', 'section']):
            score = self._calculate_content_score(element)
            if score > highest_score:
                highest_score = score
                best_candidate = element

        return best_candidate if best_candidate else (soup.body or soup)

    def _clean_html_and_extract_text(self, soup: BeautifulSoup) -> str:
        # 1단계: 불필요한 태그 일괄 제거
        for tag in soup(['script', 'style', 'header', 'footer', 'nav', 'aside', 'form', 'button', 'iframe', 'figure']):
            tag.decompose()

        # 2단계: 최고 점수 영역 선정
        best_candidate = self._get_best_candidate(soup)
        if not best_candidate: return ""

        # [최종 수정] 3단계: '분석'과 '제거'를 분리하여 안정성 확보
        # 3-1. 먼저 제거할 노이즈 태그 목록을 모두 찾습니다.
        tags_to_decompose = []
        for tag in best_candidate.find_all(['div', 'section', 'ul', 'li', 'table']):
            # 순회 중 None이 되는 경우를 원천 차단
            if not isinstance(tag, Tag): continue

            tag_id = tag.get('id', '').lower()
            tag_class = ' '.join(tag.get('class', [])).lower()
            if any(keyword in tag_id or keyword in tag_class for keyword in ['comment', 'related', 'share', 'social', 'extra', 'footer', 'header', 'nav', 'menu', 'button', 'author', 'widget']):
                tags_to_decompose.append(tag)

        # 3-2. 분석이 모두 끝난 후, 별도의 루프에서 안전하게 제거합니다.
        for tag in tags_to_decompose:
            tag.decompose()

        raw_text = best_candidate.get_text(separator='\n', strip=False)
        return self._final_text_clean(raw_text)

    def _final_text_clean(self, text: str) -> str:
        text = re.sub(r'(\n\s*){3,}', '\n\n', text)
        lines = text.split('\n')
        cleaned_lines = []
        noise_keywords = [
            '다운로드', '뷰어', '첨부파일', '목록으로', '이전글', '다음글', '맨위로',
            'Copyright', 'All rights reserved', '찾아오시는 길', '개인정보처리방침', '유용한 정보가 되었나요?'
        ]
        for line in lines:
            stripped_line = line.strip()
            if not stripped_line or len(stripped_line) < 10: continue
            if any(keyword in stripped_line for keyword in noise_keywords): continue
            if re.match(r'^\s*(작성자|등록일|조회수|담당부서|키워드|분류)\s*[:\s]', stripped_line): continue
            cleaned_lines.append(stripped_line)
        return "\n".join(cleaned_lines)

    # --- 2. 첨부 파일 처리 ---
    def _parse_pdf(self, content: bytes) -> str:
        return "\n".join([page.extract_text() for page in PdfReader(io.BytesIO(content)).pages if page.extract_text()])
    def _parse_docx(self, content: bytes) -> str:
        return "\n".join([para.text for para in Document(io.BytesIO(content)).paragraphs if para.text])
    def _parse_excel(self, content: bytes) -> str:
        xls = pd.ExcelFile(io.BytesIO(content))
        return "\n\n".join([xls.parse(sheet_name).to_string() for sheet_name in xls.sheet_names])
    def _parse_pptx(self, content: bytes) -> str:
        prs = Presentation(io.BytesIO(content))
        return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    def _parse_hwp(self, content: bytes) -> str:
        try:
            hwp = hwp5.HWPFile(io.BytesIO(content))
            text = hwp.body.text()
            if not text: raise ValueError("Extracted HWP text is empty.")
            return text
        except Exception:
            raise

    async def _save_failed_attachment(self, url: str, content_bytes: bytes, reason: str, site_identifier: str):
        try:
            site_dir = self.failed_attachments_dir / site_identifier
            site_dir.mkdir(exist_ok=True)
            timestamp = int(time.time())
            original_filename = Path(urlparse(url).path).name or f"{timestamp}.attachment"
            file_path = site_dir / f"{timestamp}_{original_filename}"
            with open(file_path, 'wb') as f: f.write(content_bytes)
            meta_path = site_dir / f"{timestamp}_{original_filename}.meta.json"
            meta_data = {
                "original_url": url, "saved_path": str(file_path), "failure_reason": reason,
                "timestamp": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime(timestamp))
            }
            with open(meta_path, 'w', encoding='utf-8') as f: json.dump(meta_data, f, ensure_ascii=False, indent=2)
            logger.warning(f"🚨 파싱 실패! 원본 파일 및 메타데이터 저장 완료: {file_path}")
        except Exception as e:
            logger.error(f"실패한 첨부 파일 저장 중 오류 발생: {e}")

    async def _extract_text_from_attachment(self, url: str, site_identifier: str) -> str:
        file_ext = Path(urlparse(url).path).suffix.lower()
        parser = self.file_parsers.get(file_ext)
        if not parser: return ""
        content_bytes = None
        try:
            async with self._session.get(url, timeout=60) as response:
                if response.status != 200: return ""
                content_bytes = await response.read()
            text = await asyncio.to_thread(parser, content_bytes)
            if not text: raise ValueError("파서가 빈 텍스트를 반환했습니다.")
            logger.info(f"✅ {file_ext.upper()} 텍스트 추출 성공: {url} ({len(text)}자)")
            return f"\n\n--- [첨부 파일 시작: {Path(url).name}] ---\n{text}\n--- [첨부 파일 끝] ---\n"
        except Exception as e:
            logger.error(f"첨부 파일({url}) 처리 실패: {e}")
            if content_bytes:
                await self._save_failed_attachment(url, content_bytes, str(e), site_identifier)
            return f"\n\n--- [첨부 파일 처리 실패: {Path(url).name}] ---\n"

    def _get_real_pdf_url_from_viewer(self, viewer_url: str, base_url: str) -> str | None:
        try:
            query_params = parse_qs(urlparse(viewer_url).query)
            if 'file' in query_params:
                return urljoin(base_url, unquote(query_params['file'][0]))
        except Exception: return None

    async def extract(self, url: str, base_url: str, site_identifier: str) -> Dict | None:
        try:
            async with self._session.get(url, timeout=30) as response:
                if response.status != 200 or 'text/html' not in response.headers.get('Content-Type', ''):
                    return None
                html = await response.text()

                soup_for_links = BeautifulSoup(html, 'html.parser')
                soup_for_text = BeautifulSoup(html, 'html.parser')

                title = soup_for_text.title.string.strip() if soup_for_text.title else url

                main_text = self._clean_html_and_extract_text(soup_for_text)

                all_links, attachment_links = [], []
                base_netloc = urlparse(base_url).netloc

                for a_tag in soup_for_links.find_all('a', href=True):
                    href = a_tag['href']
                    absolute_url = urljoin(base_url, href)

                    real_pdf_url = self._get_real_pdf_url_from_viewer(absolute_url, base_url)
                    if real_pdf_url and real_pdf_url not in attachment_links:
                        attachment_links.append(real_pdf_url); continue

                    file_ext = Path(urlparse(href).path).suffix.lower()
                    if file_ext in self.file_parsers and urlparse(absolute_url).netloc == base_netloc and absolute_url not in attachment_links:
                        attachment_links.append(absolute_url)
                    else:
                        all_links.append((absolute_url, a_tag.get_text(strip=True)))

                if attachment_links:
                    tasks = [self._extract_text_from_attachment(link, site_identifier) for link in attachment_links]
                    main_text += "".join(await asyncio.gather(*tasks))

                return {"url": url, "title": title, "main_text": main_text, "links": all_links}
        except Exception as e:
            logger.error(f"데이터 추출 중 예외 발생: {url}, {e}", exc_info=True)
            return None

    async def close_session(self):
        if self._session and not self._session.closed: await self._session.close()

